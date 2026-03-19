#!/usr/bin/env python3
"""Genera la presentazione marketing 4SEK in formato PPTX (apribile con Keynote)."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "4SEK_Presentazione.pptx")

# ── Colors ──
BG       = RGBColor(0x08, 0x0c, 0x18)
BG2      = RGBColor(0x0c, 0x12, 0x25)
SURFACE  = RGBColor(0x12, 0x1c, 0x38)
BORDER   = RGBColor(0x1e, 0x30, 0x60)
TEXT     = RGBColor(0xd8, 0xe0, 0xf0)
TEXT2    = RGBColor(0x7a, 0x8e, 0xb8)
TEXT3    = RGBColor(0x4a, 0x5a, 0x80)
BLUE     = RGBColor(0x2d, 0x7a, 0xf6)
CYAN     = RGBColor(0x00, 0xd4, 0xff)
RED      = RGBColor(0xff, 0x38, 0x44)
ORANGE   = RGBColor(0xff, 0x8c, 0x1a)
GREEN    = RGBColor(0x00, 0xe8, 0x7b)
WHITE    = RGBColor(0xff, 0xff, 0xff)

# Slide dimensions: 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.rotation = 0
    sf = shape.fill
    if fill_color:
        sf.solid()
        sf.fore_color.rgb = fill_color
    else:
        sf.background()
    ln = shape.line
    if border_color:
        ln.color.rgb = border_color
        ln.width = border_width
    else:
        ln.fill.background()
    # Minimal corner radius
    shape.adjustments[0] = 0.04
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=TEXT,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Helvetica Neue"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_label(slide, left, top, text):
    """Small cyan label like '// SECTION NAME'"""
    return add_text(slide, left, top, Inches(4), Inches(0.4),
                    f"// {text.upper()}", font_size=11, color=CYAN,
                    font_name="Menlo", bold=False)


def add_heading(slide, left, top, line1, line2="", width=Inches(8)):
    """Two-line heading: line1 bold white, line2 lighter"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = line1
    p1.font.size = Pt(40)
    p1.font.color.rgb = WHITE
    p1.font.bold = True
    p1.font.name = "Helvetica Neue"
    if line2:
        p2 = tf.add_paragraph()
        p2.text = line2
        p2.font.size = Pt(40)
        p2.font.color.rgb = CYAN
        p2.font.bold = False
        p2.font.name = "Georgia"
    return txBox


def add_sub(slide, left, top, text, width=Inches(6)):
    return add_text(slide, left, top, width, Inches(1),
                    text, font_size=16, color=TEXT2)


def add_image_safe(slide, path, left, top, width=None, height=None):
    full = os.path.join(BASE, path)
    if not os.path.exists(full):
        return None
    if width and height:
        return slide.shapes.add_picture(full, left, top, width, height)
    elif width:
        return slide.shapes.add_picture(full, left, top, width=width)
    elif height:
        return slide.shapes.add_picture(full, left, top, height=height)
    return slide.shapes.add_picture(full, left, top)


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    return shape


def card_with_icon(slide, left, top, width, height, icon_text, title, desc, accent=BLUE):
    """A feature card with icon, title, description"""
    card = add_rect(slide, left, top, width, height, fill_color=SURFACE, border_color=BORDER, border_width=Pt(1))
    # Accent line top
    add_rect(slide, left + Emu(20000), top, Inches(0.6), Pt(3), fill_color=accent)
    # Icon
    icon_box = add_rect(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.55), Inches(0.55), fill_color=BG2)
    add_text(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.55), Inches(0.55),
             icon_text, font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_text(slide, left + Inches(0.3), top + Inches(1.0), width - Inches(0.6), Inches(0.4),
             title, font_size=15, color=WHITE, bold=True)
    # Desc
    add_text(slide, left + Inches(0.3), top + Inches(1.4), width - Inches(0.6), Inches(1.2),
             desc, font_size=11, color=TEXT2)


# ═══════════════════════════════════════════════════════════
# SLIDE 1: TITLE / HERO
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG)

# Decorative gradient bar top
add_rect(s, Inches(0), Inches(0), SLIDE_W, Pt(3), fill_color=BLUE)

# Logo area
add_rect(s, Inches(0.8), Inches(0.6), Inches(0.65), Inches(0.65), fill_color=BLUE)
add_text(s, Inches(0.8), Inches(0.63), Inches(0.65), Inches(0.65),
         "4S", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Menlo")
add_text(s, Inches(1.6), Inches(0.65), Inches(2), Inches(0.5),
         "4SEK", font_size=22, color=WHITE, bold=True, font_name="Menlo")

# Badge
add_rect(s, Inches(0.8), Inches(2.0), Inches(2.6), Inches(0.35), fill_color=BG2, border_color=BORDER, border_width=Pt(1))
add_text(s, Inches(0.9), Inches(2.02), Inches(2.4), Inches(0.35),
         "● PIATTAFORMA AI-POWERED", font_size=9, color=CYAN, font_name="Menlo")

# Main title
txBox = s.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(6.5), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p1 = tf.paragraphs[0]
p1.text = "Sicurezza"
p1.font.size = Pt(60)
p1.font.color.rgb = BLUE
p1.font.bold = True
p1.font.name = "Helvetica Neue"
p2 = tf.add_paragraph()
p2.text = "Intelligente"
p2.font.size = Pt(60)
p2.font.color.rgb = TEXT
p2.font.bold = False
p2.font.name = "Georgia"
p3 = tf.add_paragraph()
p3.text = "per Centri Commerciali"
p3.font.size = Pt(44)
p3.font.color.rgb = TEXT
p3.font.bold = True
p3.font.name = "Helvetica Neue"

# Subtitle
add_text(s, Inches(0.8), Inches(5.2), Inches(5.5), Inches(0.9),
         "Telecamere con AI a bordo, rilevamento minacce in tempo reale, "
         "coordinamento guardianie e risposta automatica agli eventi critici.",
         font_size=16, color=TEXT2)

# Camera image
add_image_safe(s, "assets/img/smart_camera.png",
               Inches(8.5), Inches(1.5), width=Inches(4))

# Floating tags around camera
add_rect(s, Inches(8.2), Inches(2.0), Inches(2.2), Inches(0.35), fill_color=SURFACE, border_color=BORDER, border_width=Pt(1))
add_text(s, Inches(8.3), Inches(2.02), Inches(2), Inches(0.35),
         "AI ON-BOARD • Edge", font_size=9, color=CYAN, font_name="Menlo")

add_rect(s, Inches(10.5), Inches(3.8), Inches(2.2), Inches(0.35), fill_color=SURFACE, border_color=BORDER, border_width=Pt(1))
add_text(s, Inches(10.6), Inches(3.82), Inches(2), Inches(0.35),
         "Confidence: 96.3%", font_size=9, color=CYAN, font_name="Menlo")

add_rect(s, Inches(8.0), Inches(5.0), Inches(2.4), Inches(0.35), fill_color=BG2, border_color=RGBColor(0x40, 0x18, 0x18), border_width=Pt(1))
add_text(s, Inches(8.1), Inches(5.02), Inches(2.2), Inches(0.35),
         "⚠ THREAT DETECTED", font_size=9, color=RED, font_name="Menlo")

# Bottom bar
add_rect(s, Inches(0), Inches(7.15), SLIDE_W, Pt(1), fill_color=BORDER)
add_text(s, Inches(0.8), Inches(7.0), Inches(4), Inches(0.4),
         "4SEK — Sicurezza Intelligente", font_size=10, color=TEXT3, font_name="Menlo")


# ═══════════════════════════════════════════════════════════
# SLIDE 2: IL PROBLEMA
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Pt(3), fill_color=RED)

add_label(s, Inches(0.8), Inches(0.8), "Il problema")
add_heading(s, Inches(0.8), Inches(1.3), "La sicurezza tradizionale", "non basta più")
add_sub(s, Inches(0.8), Inches(2.8),
        "I sistemi convenzionali si basano sulla sorveglianza passiva. "
        "Gli operatori non possono monitorare tutto simultaneamente.")

# 3 problem cards
card_w = Inches(3.7)
card_h = Inches(3.2)
gap = Inches(0.3)
start_x = Inches(0.8)
y = Inches(3.8)

problems = [
    ("👁", "Sorveglianza passiva",
     "Gli operatori CCTV possono monitorare efficacemente solo 4-6 schermi. "
     "Il resto è un punto cieco. Gli eventi critici vengono identificati solo a posteriori.", RED),
    ("⏱", "Tempi di reazione lenti",
     "Tra il rilevamento di un evento e l'intervento delle guardie passano in media 4-8 minuti. "
     "Troppo per prevenire scippi e aggressioni.", ORANGE),
    ("📊", "Zero prevenzione",
     "Senza analisi predittiva, i sistemi tradizionali reagiscono solo dopo l'evento. "
     "Nessuna capacità di identificare comportamenti sospetti in anticipo.", BLUE),
]
for i, (icon, title, desc, accent) in enumerate(problems):
    x = start_x + i * (card_w + gap)
    card_with_icon(s, x, y, card_w, card_h, icon, title, desc, accent)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: AI VISION (1/2) — Rilevamento & Heatmap
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Pt(3), fill_color=CYAN)

add_label(s, Inches(0.8), Inches(0.6), "Visione artificiale")
add_heading(s, Inches(0.8), Inches(1.0), "L'AI che vede ciò che", "l'occhio umano perde")

# Image 1: Assault detection
img1 = add_image_safe(s, "img_ai/Screenshot 2026-03-18 at 09.11.26.png",
                       Inches(0.8), Inches(3.0), width=Inches(5.8))

# Label for image 1
add_rect(s, Inches(0.8), Inches(5.8), Inches(2.6), Inches(0.35), fill_color=BG2, border_color=RGBColor(0x40, 0x15, 0x15), border_width=Pt(1))
add_text(s, Inches(0.9), Inches(5.82), Inches(2.4), Inches(0.35),
         "● RILEVAMENTO ATTIVO", font_size=9, color=RED, font_name="Menlo")
add_text(s, Inches(0.8), Inches(6.2), Inches(5.5), Inches(0.5),
         "Rilevamento Aggressioni — Bounding box in tempo reale con tracking del soggetto, velocità e direzione di fuga",
         font_size=13, color=TEXT2)

# Image 2: Heatmap
img2 = add_image_safe(s, "img_ai/Screenshot 2026-03-18 at 09.11.34.png",
                       Inches(7.0), Inches(3.0), width=Inches(5.8))

add_rect(s, Inches(7.0), Inches(5.8), Inches(2.2), Inches(0.35), fill_color=BG2, border_color=RGBColor(0x10, 0x30, 0x40), border_width=Pt(1))
add_text(s, Inches(7.1), Inches(5.82), Inches(2), Inches(0.35),
         "● ANALISI FLUSSI", font_size=9, color=CYAN, font_name="Menlo")
add_text(s, Inches(7.0), Inches(6.2), Inches(5.5), Inches(0.5),
         "Heatmap & Assembramenti — Densità real-time, conteggio per zona, allerta automatica al superamento soglia",
         font_size=13, color=TEXT2)


# ═══════════════════════════════════════════════════════════
# SLIDE 4: AI VISION (2/2) — Parcheggi
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Pt(3), fill_color=CYAN)

add_label(s, Inches(0.8), Inches(0.6), "Visione artificiale")
add_heading(s, Inches(0.8), Inches(1.0), "Sorveglianza parcheggi", "con AI integrata")

# Image 3: Tampering
img3 = add_image_safe(s, "img_ai/Screenshot 2026-03-18 at 09.10.48.png",
                       Inches(0.8), Inches(3.0), width=Inches(5.8))

add_rect(s, Inches(0.8), Inches(5.5), Inches(2.8), Inches(0.35), fill_color=BG2, border_color=RGBColor(0x40, 0x28, 0x08), border_width=Pt(1))
add_text(s, Inches(0.9), Inches(5.52), Inches(2.6), Inches(0.35),
         "⚠ ALERT MANOMISSIONE", font_size=9, color=ORANGE, font_name="Menlo")
add_text(s, Inches(0.8), Inches(5.9), Inches(5.5), Inches(0.7),
         "Rilevamento Manomissioni — Identificazione tentativi di effrazione veicoli, "
         "soggetti sospetti, prossimità anomala",
         font_size=13, color=TEXT2)

# Image 4: Parking overview
img4 = add_image_safe(s, "img_ai/Screenshot 2026-03-18 at 09.10.33.png",
                       Inches(7.0), Inches(3.0), width=Inches(5.8))

add_rect(s, Inches(7.0), Inches(5.5), Inches(2.8), Inches(0.35), fill_color=BG2, border_color=RGBColor(0x10, 0x20, 0x40), border_width=Pt(1))
add_text(s, Inches(7.1), Inches(5.52), Inches(2.6), Inches(0.35),
         "● MONITORAGGIO 24/7", font_size=9, color=BLUE, font_name="Menlo")
add_text(s, Inches(7.0), Inches(5.9), Inches(5.5), Inches(0.7),
         "Sorveglianza Parcheggi — Copertura multilivello con telecamere a LED di stato integrato, "
         "monitoraggio continuo di ogni corsia",
         font_size=13, color=TEXT2)


# ═══════════════════════════════════════════════════════════
# SLIDE 5: FUNZIONALITÀ
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Pt(3), fill_color=BLUE)

add_label(s, Inches(0.8), Inches(0.5), "Funzionalità")
add_heading(s, Inches(0.8), Inches(0.9), "Una piattaforma,", "protezione totale")

features = [
    ("⚠", "Rilevamento Minacce\nReal-time",
     "Armi, risse, furti e scippi identificati in millisecondi con modelli neurali on-board.", RED),
    ("🎤", "Analisi Audio\nIntelligente",
     "Microfoni integrati rilevano urla, richieste d'aiuto e suoni anomali in automatico.", ORANGE),
    ("🌡", "Heatmap &\nAnalisi Flussi",
     "Densità persone in real-time, rilevamento assembramenti e stazionamento prolungato.", CYAN),
    ("🔊", "Altoparlanti\nIntelligenti",
     "Messaggi deterrenti automatici, allarmi e indicazioni di evacuazione in tempo reale.", GREEN),
    ("🛡", "Coordinamento\nGuardianie",
     "Dispatch automatico, GPS tracking pattuglie, gestione ronde con checkpoint.", BLUE),
    ("🚨", "Integrazione\nForze dell'Ordine",
     "Chiamate auto a 112/118/115, condivisione video real-time con le autorità.", RGBColor(0xa8, 0x55, 0xf7)),
]

card_w = Inches(3.7)
card_h = Inches(2.8)
gap_x = Inches(0.3)
gap_y = Inches(0.25)
start_x = Inches(0.8)
start_y = Inches(2.6)

for i, (icon, title, desc, accent) in enumerate(features):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    card_with_icon(s, x, y, card_w, card_h, icon, title, desc, accent)


# ═══════════════════════════════════════════════════════════
# SLIDE 6: CENTRALE OPERATIVA (screenshot)
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Pt(3), fill_color=BLUE)

add_label(s, Inches(0.8), Inches(0.5), "Centrale Operativa")
add_heading(s, Inches(0.8), Inches(0.9), "Controllo totale", "da un'unica interfaccia")
add_sub(s, Inches(0.8), Inches(2.3),
        "Dashboard operativa con vista live su tutte le telecamere, "
        "gestione allarmi, mappa interattiva, analitiche e comunicazioni.",
        width=Inches(7))

# Monitor frame
monitor_x = Inches(0.8)
monitor_y = Inches(3.2)
monitor_w = Inches(11.7)
monitor_h = Inches(3.8)

# Bezel
add_rect(s, monitor_x, monitor_y, monitor_w, monitor_h,
         fill_color=RGBColor(0x0f, 0x0f, 0x14), border_color=RGBColor(0x1a, 0x1a, 0x24), border_width=Pt(2))

# Webcam dot
add_rect(s, Inches(6.55), monitor_y + Emu(20000), Inches(0.12), Inches(0.12),
         fill_color=RGBColor(0x1a, 0x1a, 0x22))

# Screen area with screenshot - use one of the gallery images as placeholder for the dashboard
# We'll use the existing images to compose a representative view
screen_pad = Inches(0.15)
screen_x = monitor_x + screen_pad
screen_y = monitor_y + Inches(0.25)
screen_w = monitor_w - screen_pad * 2
screen_h = monitor_h - Inches(0.35)

# Dark screen background
add_rect(s, screen_x, screen_y, screen_w, screen_h, fill_color=RGBColor(0x0a, 0x0e, 0x17))

# Place a representative camera image inside the "screen"
add_image_safe(s, "assets/img/gallery_heatmap.png",
               screen_x + Inches(0.15), screen_y + Inches(0.15), width=Inches(3.8))
add_image_safe(s, "assets/img/plaza_detection.png",
               screen_x + Inches(4.1), screen_y + Inches(0.15), width=Inches(3.8))
add_image_safe(s, "assets/img/parking_detection.png",
               screen_x + Inches(8.0), screen_y + Inches(0.15), width=Inches(3.2))

# Labels inside screen
add_text(s, screen_x + Inches(0.2), screen_y + screen_h - Inches(0.5), Inches(3), Inches(0.3),
         "CAM-INT-03 • Atrio Centrale", font_size=9, color=CYAN, font_name="Menlo")
add_text(s, screen_x + Inches(4.2), screen_y + screen_h - Inches(0.5), Inches(3), Inches(0.3),
         "CAM-EXT-07 • Piazza Sud", font_size=9, color=RED, font_name="Menlo")
add_text(s, screen_x + Inches(8.1), screen_y + screen_h - Inches(0.5), Inches(3), Inches(0.3),
         "CAM-PK-02 • Parcheggio B2", font_size=9, color=ORANGE, font_name="Menlo")

# Stand
stand_x = Inches(5.6)
add_rect(s, stand_x, monitor_y + monitor_h, Inches(2), Inches(0.3),
         fill_color=RGBColor(0x0f, 0x0f, 0x14))

# Annotations
annots = [
    (Inches(0.1), Inches(3.7), "📹 Live camera grid"),
    (Inches(12.0), Inches(4.5), "⚠ Alert panel real-time"),
    (Inches(0.1), Inches(6.0), "📊 Metriche & KPI"),
]
for ax, ay, atxt in annots:
    add_rect(s, ax, ay, Inches(2.3), Inches(0.35), fill_color=BG2, border_color=RGBColor(0x00, 0x40, 0x50), border_width=Pt(1))
    add_text(s, ax + Inches(0.1), ay + Inches(0.02), Inches(2.1), Inches(0.35),
             atxt, font_size=9, color=CYAN, font_name="Menlo")


# ═══════════════════════════════════════════════════════════
# SLIDE 7: APP MOBILE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Pt(3), fill_color=GREEN)

add_label(s, Inches(0.8), Inches(0.5), "App mobile")
add_heading(s, Inches(0.8), Inches(0.9), "4SEK Guard", "nelle mani del tuo team")

# Phone frame
phone_x = Inches(8.8)
phone_y = Inches(1.0)
phone_w = Inches(3.2)
phone_h = Inches(6.2)

# Phone bezel
add_rect(s, phone_x, phone_y, phone_w, phone_h,
         fill_color=RGBColor(0x0a, 0x0a, 0x10), border_color=RGBColor(0x1e, 0x1e, 0x2a), border_width=Pt(2))

# Notch
add_rect(s, phone_x + Inches(0.8), phone_y, Inches(1.6), Inches(0.3),
         fill_color=RGBColor(0x0a, 0x0a, 0x10))

# Screen content using a relevant image
screen_px = phone_x + Inches(0.15)
screen_py = phone_y + Inches(0.35)
screen_pw = phone_w - Inches(0.3)
screen_ph = phone_h - Inches(0.5)
add_rect(s, screen_px, screen_py, screen_pw, screen_ph, fill_color=RGBColor(0x05, 0x08, 0x0f))

# Place camera images inside phone to simulate the app
add_image_safe(s, "assets/img/plaza_detection.png",
               screen_px + Inches(0.1), screen_py + Inches(1.3), width=screen_pw - Inches(0.2))

# App header text
add_text(s, screen_px + Inches(0.15), screen_py + Inches(0.15), Inches(2), Inches(0.35),
         "Dashboard", font_size=14, color=WHITE, bold=True)
add_text(s, screen_px + Inches(0.15), screen_py + Inches(0.5), Inches(2.2), Inches(0.25),
         "Turno 06:00 - 14:00 • G3 Verdi", font_size=7, color=TEXT3, font_name="Menlo")

# Alert banner in phone
add_rect(s, screen_px + Inches(0.1), screen_py + Inches(0.85), screen_pw - Inches(0.2), Inches(0.4),
         fill_color=RGBColor(0x20, 0x08, 0x08), border_color=RGBColor(0x50, 0x15, 0x15), border_width=Pt(1))
add_text(s, screen_px + Inches(0.2), screen_py + Inches(0.87), Inches(2), Inches(0.2),
         "⚠ CRITICO", font_size=7, color=RED, bold=True, font_name="Menlo")
add_text(s, screen_px + Inches(0.2), screen_py + Inches(1.05), Inches(2.2), Inches(0.2),
         "Aggressione - Piazza Sud", font_size=8, color=WHITE, bold=True)

# Features list on the left
mobile_features = [
    ("📱", "Dashboard Mobile",
     "Panoramica istantanea: allerte attive, statistiche in tempo reale, stato del team."),
    ("⚠", "Dettaglio Allarmi",
     "Feed camera live, mappa con navigazione all'incidente, dati AI completi."),
    ("📡", "Radio & Comunicazioni",
     "Push-to-talk integrato, messaggi real-time dalla centrale, contatti rapidi."),
    ("📋", "Report & Ronde",
     "Gestione ronde con checkpoint, report fotografici, storico interventi."),
]

fy = Inches(2.5)
for icon, title, desc in mobile_features:
    # Icon box
    add_rect(s, Inches(0.8), fy, Inches(0.5), Inches(0.5), fill_color=SURFACE, border_color=BORDER, border_width=Pt(1))
    add_text(s, Inches(0.8), fy + Inches(0.02), Inches(0.5), Inches(0.5),
             icon, font_size=16, alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), fy, Inches(5), Inches(0.3),
             title, font_size=15, color=WHITE, bold=True)
    add_text(s, Inches(1.5), fy + Inches(0.35), Inches(5.5), Inches(0.5),
             desc, font_size=12, color=TEXT2)
    fy += Inches(1.1)


# ═══════════════════════════════════════════════════════════
# SLIDE 8: STATISTICHE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Pt(3), fill_color=CYAN)

add_label(s, Inches(0.8), Inches(0.8), "Numeri")
add_heading(s, Inches(0.8), Inches(1.3), "I numeri parlano", "da soli")

stats = [
    ("30+", "Telecamere AI", "Edge computing on-board"),
    ("<3s", "Tempo di Reazione", "Dal rilevamento all'alert"),
    ("96.3%", "Accuratezza AI", "Confidence media detection"),
    ("24/7", "Monitoraggio", "Zero interruzioni"),
]

card_w = Inches(2.7)
card_h = Inches(3.0)
gap = Inches(0.3)
total_w = 4 * card_w + 3 * gap
start_x = (SLIDE_W - total_w) / 2
y = Inches(3.2)

for i, (val, label, sub) in enumerate(stats):
    x = start_x + i * (card_w + gap)
    card = add_rect(s, x, y, card_w, card_h, fill_color=SURFACE, border_color=BORDER, border_width=Pt(1))
    # Top accent
    add_rect(s, x + Inches(0.6), y, Inches(1.5), Pt(2), fill_color=CYAN)
    # Value
    add_text(s, x, y + Inches(0.6), card_w, Inches(0.8),
             val, font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Menlo")
    # Label
    add_text(s, x, y + Inches(1.5), card_w, Inches(0.4),
             label, font_size=15, color=TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    # Sub
    add_text(s, x, y + Inches(2.0), card_w, Inches(0.4),
             sub, font_size=11, color=TEXT3, alignment=PP_ALIGN.CENTER, font_name="Menlo")


# ═══════════════════════════════════════════════════════════
# SLIDE 9: CTA
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Pt(3), fill_color=BLUE)

# Centered content
add_text(s, Inches(0), Inches(1.5), SLIDE_W, Inches(0.4),
         "// INIZIA ORA", font_size=11, color=CYAN, font_name="Menlo", alignment=PP_ALIGN.CENTER)

txBox = s.shapes.add_textbox(Inches(2), Inches(2.2), Inches(9.3), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p1 = tf.paragraphs[0]
p1.text = "Proteggi il tuo"
p1.font.size = Pt(48)
p1.font.color.rgb = WHITE
p1.font.bold = True
p1.font.name = "Helvetica Neue"
p1.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "centro commerciale"
p2.font.size = Pt(48)
p2.font.color.rgb = CYAN
p2.font.bold = False
p2.font.name = "Georgia"
p2.alignment = PP_ALIGN.CENTER

add_text(s, Inches(3), Inches(4.2), Inches(7.3), Inches(0.8),
         "Richiedi una demo personalizzata e scopri come 4SEK\n"
         "trasforma la sicurezza del tuo spazio commerciale.",
         font_size=16, color=TEXT2, alignment=PP_ALIGN.CENTER)

# CTA Button
btn = add_rect(s, Inches(4.8), Inches(5.3), Inches(3.7), Inches(0.7),
               fill_color=BLUE)
add_text(s, Inches(4.8), Inches(5.35), Inches(3.7), Inches(0.65),
         "Richiedi una Demo", font_size=18, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

# Trust badges
add_text(s, Inches(0), Inches(6.3), SLIDE_W, Inches(0.4),
         "✓ Setup in 48 ore    ✓ Nessun costo nascosto    ✓ Supporto 24/7",
         font_size=12, color=TEXT3, alignment=PP_ALIGN.CENTER, font_name="Menlo")

# Footer
add_rect(s, Inches(0), Inches(7.15), SLIDE_W, Pt(1), fill_color=BORDER)
add_text(s, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3),
         "© 2026 4SEK — Sicurezza Intelligente", font_size=10, color=TEXT3,
         alignment=PP_ALIGN.CENTER, font_name="Menlo")


# ═══ SAVE ═══
prs.save(OUT)
print(f"✓ Presentazione salvata: {OUT}")
